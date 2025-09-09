import io, zipfile, csv, random
import streamlit as st
# --- helper: pick target words from a sentence ---
import re
from pathlib import Path



def load_css(path: str | Path) -> None:
    css_path = Path(path)
    css = css_path.read_text(encoding="utf-8")
    st.markdown(f"<style>{css}</style>", unsafe_allow_html=True)





# ---- preview helper: scale an SVG to a target pixel width (keeps aspect) ----

def _scale_svg_for_preview(svg_text: str, target_width_px: int) -> tuple[str, int]:
    """
    Returns (scaled_svg_text, new_height_px).
    Only used for UI preview; original SVGs stay full size for ZIP/PNG/PDF.
    """
    s = svg_text
    # Prefer viewBox; if missing, add one from width/height.
    m = re.search(r'viewBox="0\s+0\s+([\d.]+)\s+([\d.]+)"', s)
    if m:
        vw, vh = float(m.group(1)), float(m.group(2))
    else:
        mw = re.search(r'(<svg\b[^>]*\bwidth=")([^"]+)(")', s)
        mh = re.search(r'(<svg\b[^>]*\bheight=")([^"]+)(")', s)
        if not (mw and mh):
            return s, 600  # fallback
        try:
            vw = float(re.sub(r'[^0-9.]+', '', mw.group(2)))
            vh = float(re.sub(r'[^0-9.]+', '', mh.group(2)))
        except Exception:
            return s, 600
        # insert a viewBox
        s = re.sub(r'<svg\b', f'<svg viewBox="0 0 {vw:.2f} {vh:.2f}"', s, count=1)

    scale = max(0.05, float(target_width_px) / max(1.0, vw))
    new_h = max(50, int(round(vh * scale)))

    # rewrite width/height only on the <svg ...> tag
    s = re.sub(r'(<svg\b[^>]*\bwidth=")[^"]+(")',  rf'\g<1>{int(target_width_px)}\g<2>', s, count=1)
    s = re.sub(r'(<svg\b[^>]*\bheight=")[^"]+(")', rf'\g<1>{new_h}\g<2>',            s, count=1)
    # ensure preserve ratio (optional but nice)
    if 'preserveAspectRatio' not in s[:400]:
        s = re.sub(r'<svg\b', '<svg preserveAspectRatio="xMidYMid meet"', s, count=1)
    return s, new_h

def _targets_from_sentence(sentence: str, n: int, mn: int, mx: int, seed: str | None = None):
    """
    Very simple tokenizer -> unique UPPER words within [mn, mx] length,
    shuffled (deterministically if seed given), then take first n.
    """
    toks = re.findall(r"[A-Za-z0-9]+", sentence or "")
    seen = set()
    uniq = []
    for t in toks:
        up = t.upper()
        if mn <= len(up) <= mx and up not in seen:
            seen.add(up)
            uniq.append(up)
    if not uniq:
        return []
    rng = random.Random(seed or None)
    rng.shuffle(uniq)
    return uniq[: max(0, int(n or 0))]


st.set_page_config(page_title="Puzzle Generator", layout="wide")
# If styles.css is next to app.py:
load_css(Path(__file__).with_name("styles.css"))
st.title("Puzzle Generator (Web – Simple)")




# --- Controls in the sidebar (clean + compact) ---
with st.sidebar:
    # st.header("Controls")

    # Real tabs inside the sidebar
    tab_create, tab_settings = st.tabs(["Create Puzzle", "Settings"])

    # ---------------------------
    # TAB 1: Create Puzzle
    # ---------------------------
    with tab_create:
        gen = st.selectbox("Puzzle type", ["Word Search – Sentence"])

        # Row 1: Grid width | Grid height
        r1c1, r1c2 = st.columns(2)
        with r1c1:
            grid_w = st.number_input("Grid width", 5, 40, 15, format="%d")
        with r1c2:
            grid_h = st.number_input("Grid height", 5, 40, 15, format="%d")

        # Row 2: # puzzles | # words per puzzle
        r2c1, r2c2 = st.columns(2)
        with r2c1:
            n_puzzles = st.number_input("# puzzles", 1, 100, 6, format="%d")
        with r2c2:
            num_words = st.number_input("# words per puzzle", 1, 40, 10, format="%d")

        # Row 3: Min | Max word length
        r3c1, r3c2 = st.columns(2)
        with r3c1:
            min_len = st.number_input("Min word length", 2, 20, 3, format="%d")
        with r3c2:
            max_len = st.number_input("Max word length", 2, 20, 10, format="%d")

        # st.divider()
        seed = st.text_input("Seed (optional)", "")

        # Uploader spans full sidebar width (not inside a sub-column)
        csv_file = st.file_uploader("CSV: Title, Sentence, Secret", type=["csv"])

        # st.divider()
        # size_label = st.select_slider("Preview size", options=["Small","Medium","Large"], value="Medium")
        # PREVIEW_W = {"Small": 420, "Medium": 560, "Large": 720}[size_label]

        go = st.button("Generate", type="primary", use_container_width=True, disabled=(csv_file is None))

    # ---------------------------
    # TAB 2: Settings
    # ---------------------------
    with tab_settings:
        st.caption("Output formats")
        make_png  = st.checkbox("Also make PNG", value=True)
        make_pdf  = st.checkbox("Also make PDF", value=False)
        make_pptx = st.checkbox("Also make PPTX (simple insert)", value=False)

        # st.divider()
        st.caption("Preview")
        size_label2 = st.select_slider("Preview size", options=["Small","Medium","Large"], value="Medium")
        PREVIEW_W = {"Small": 420, "Medium": 560, "Large": 720}[size_label2]





if go:
    # --- Import inside the button, so errors show on page ---
    try:
        import puzzle_engine as eng
    except Exception as e:
        st.error("Failed to import puzzle_engine.py")
        st.exception(e)
        st.stop()

    try:
        import svg_renderer as svg
    except Exception as e:
        st.error("Failed to import svg_renderer.py")
        st.exception(e)
        st.stop()

    try:
        from cairosvg import svg2png, svg2pdf
    except Exception as e:
        st.error("cairosvg not installed or failed to import")
        st.exception(e)
        st.stop()

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        if make_pptx:
            st.error("python-pptx failed to import")
            st.exception(e)
            st.stop()
        else:
            Presentation = None  # not used

    # --- Read CSV ---
    rows = []
    try:
        for i, r in enumerate(csv.reader(io.TextIOWrapper(csv_file, encoding="utf-8-sig"))):
            if i == 0: 
                continue
            title = (r[0] if len(r)>0 else "").strip()
            sentence = (r[1] if len(r)>1 else "").strip()
            secret = (r[2] if len(r)>2 else "").strip()
            if sentence:
                rows.append((title, sentence, secret))
    except Exception as e:
        st.error("Could not read CSV")
        st.exception(e)
        st.stop()

    if not rows:
        st.error("No sentence rows found.")
        st.stop()

    rng = random.Random(seed or None)
    rng.shuffle(rows)
    picks = rows[:n_puzzles]

    svgs = []
    imgs_for_pptx = []
    first_puz_svg = None
    first_sol_svg = None


    try:
        for idx, (title, sentence, secret) in enumerate(picks, 1):
            # 1) derive target words from the sentence
            targets = _targets_from_sentence(
                sentence=sentence,
                n=num_words,
                mn=min_len,
                mx=max_len,
                seed=seed or None,
            )

            if not targets:
                st.warning(f"No usable words found in sentence #{idx}; skipping.")
                continue

            # 2) build spec — IMPORTANT: pass words=targets and sync num_words
            spec = eng.PuzzleSpec(
                mode="Sentence",  # label
                grid_width=grid_w, grid_height=grid_h,
                num_words=len(targets),
                min_len=min_len, max_len=max_len,
                order_policy="source",          # keep the order we chose
                runout_policy="stop",
                directions=["N","NE","E","SE","S","SW","W","NW"],
                seed=seed or None,
                words=targets,                  # <<< this makes placement happen
                sentences=[sentence],
                hidden_text=(secret or None),
            )

            # 3) generate and render with sentence + targets for legend & underline
            res = eng.generate_one_puzzle(spec)

            look = svg.Appearance(
                grid_font_family="Arial",
                grid_font_size=24,
                legend_columns=1,              # irrelevant now, but harmless
                show_legend=False,             # <<< HIDE word-list legend on the PUZZLE
                # --- sentence block settings ---
                show_sentence=True,            # <<< draw sentence block in PUZZLE
                sentence_text=sentence,        # full sentence
                sentence_targets=targets,      # words to underline/bold in sentence
                sentence_font_family="Arial",
                sentence_font_size=16,
                sentence_font_bold=False,
                sentence_font_color="#000000",
                sentence_align="Left",
                sentence_width_factor=0.459,
                # --- solution settings ---
                solution_show_legend=False,    # keep solution clean
                solution_show_sentence=bool(secret),  # show sentence also in SOLUTION if you like
                solution_mark_style="highlight",
                solution_mark_color="#D94242",
            )

            puz_svg = svg.render_puzzle_svg(res, look)
            # inject the sentence block into the PUZZLE svg
            puz_svg = svg.inject_sentence_block(puz_svg, look)

            sol_svg = svg.render_solution_svg(res, look)
            # inject the sentence block into the SOLUTION svg only if requested
            if look.solution_show_sentence:
                sol_svg = svg.inject_sentence_block(sol_svg, look)

            svgs.append((f"puzzle_{idx:03d}.svg", puz_svg))
            svgs.append((f"solution_{idx:03d}.svg", sol_svg))

            if first_puz_svg is None:
                first_puz_svg = puz_svg
                first_sol_svg = sol_svg


    except Exception as e:
        st.error("Puzzle generation/rendering failed")
        st.exception(e)
        st.stop()

   
    # --- Previews (tabs) ---
    tab_puz, tab_sol = st.tabs(["Preview — Puzzle", "Preview — Solution"])

    with tab_puz:
        if first_puz_svg:
            svgp, hp = _scale_svg_for_preview(first_puz_svg, PREVIEW_W)
            st.components.v1.html(svgp, height=hp + 6, scrolling=False)
        else:
            st.info("No preview available.")

    with tab_sol:
        if first_sol_svg:
            svg_sol_preview, hs = _scale_svg_for_preview(first_sol_svg, PREVIEW_W)
            st.components.v1.html(svg_sol_preview, height=hs + 6, scrolling=False)
        else:
            st.info("No preview available.")







    # --- ZIP outputs ---
    try:
        mem = io.BytesIO()
        with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, s in svgs:
                zf.writestr(name, s)

            if make_png or make_pdf or make_pptx:
                from cairosvg import svg2png, svg2pdf
                for name, s in svgs:
                    try:
                        if make_png:
                            zf.writestr(name.replace(".svg", ".png"),
                                        svg2png(bytestring=s.encode("utf-8")))
                    except Exception as e:
                        zf.writestr(name.replace(".svg", ".PNG_ERROR.txt"),
                                    (f"PNG conversion failed for {name}:\n{e}").encode("utf-8"))

                    try:
                        if make_pdf:
                            zf.writestr(name.replace(".svg", ".pdf"),
                                        svg2pdf(bytestring=s.encode("utf-8")))
                    except Exception as e:
                        zf.writestr(name.replace(".svg", ".PDF_ERROR.txt"),
                                    (f"PDF conversion failed for {name}:\n{e}").encode("utf-8"))

                    try:
                        if make_pptx and name.startswith("puzzle_"):
                            imgs_for_pptx.append(svg2png(bytestring=s.encode("utf-8")))
                    except Exception as e:
                        zf.writestr(name.replace(".svg", ".PPTX_IMAGE_ERROR.txt"),
                                    (f"PPTX image prep failed for {name}:\n{e}").encode("utf-8"))


            if make_pptx and imgs_for_pptx:
                prs = Presentation()
                blank = prs.slide_layouts[6]
                import io as _io
                for png in imgs_for_pptx:
                    slide = prs.slides.add_slide(blank)
                    stream = _io.BytesIO(png)
                    slide.shapes.add_picture(stream, Inches(0.5), Inches(0.5), height=Inches(7.5))
                out = _io.BytesIO(); prs.save(out)
                zf.writestr("puzzles.pptx", out.getvalue())

        mem.seek(0)
        st.download_button("Download ZIP", data=mem.read(), file_name="puzzles.zip", mime="application/zip")
    except Exception as e:
        st.error("Failed to package outputs")
        st.exception(e)
        st.stop()
